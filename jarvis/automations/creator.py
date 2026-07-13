"""File-creator automations — text formats, binary documents, images, and code execution.

Tools registered:
    create_markdown, create_html, create_csv, create_json, create_yaml,
    create_xml, create_txt, create_python_file, create_code_file,
    create_pdf, create_docx, create_xlsx, create_pptx,
    generate_image, run_python.

Design notes
------------
- All file outputs default to ~/jarvis_workspace if no path is given, but a
  full OS path also works (we just deny writes into protected system dirs).
- For PDF / DOCX / XLSX / PPTX, the LLM only has to write Markdown-lite:
  `#` / `##` / `###` headings, blank-separated paragraphs, `-` bullets,
  `>` quotes, and `| col | col |` tables. `_md_to_blocks` parses it and each
  format renders it the way it likes. This keeps the prompt tiny.
- Image generation hits the gemini-2.5-flash-image model with the same
  key-rotation / fallback as the Brain.
- run_python is approval-gated like run_command — a 10s timeout in a temp
  working dir is the safety net. Not a hermetic sandbox.
"""
from __future__ import annotations

import base64
import csv as _csv
import io
import json
import re
import subprocess
import sys
import tempfile
import time
from datetime import datetime
from html import escape as _html_escape
from pathlib import Path
from xml.sax.saxutils import escape as _xml_escape

import requests

from config import CONFIG, GEMINI_BASE

from .registry import tool


# ── shared helpers ───────────────────────────────────────
WORKSPACE = Path.home() / "jarvis_workspace"
WORKSPACE.mkdir(parents=True, exist_ok=True)
IMAGES_DIR = WORKSPACE / "images"
IMAGES_DIR.mkdir(parents=True, exist_ok=True)

# Same denylist as files.py — refuse to write inside protected system dirs.
_PROTECTED = [
    Path("C:/Windows"), Path("C:/Windows/System32"),
    Path("C:/Program Files"), Path("C:/Program Files (x86)"),
    Path("C:/ProgramData"),
    Path("/System"), Path("/usr"), Path("/bin"), Path("/sbin"), Path("/etc"),
]


def _is_blocked(path: str) -> bool:
    try:
        p = Path(path).expanduser().resolve()
    except OSError:
        return True
    for b in _PROTECTED:
        try:
            if p == b or b in p.parents:
                return True
        except OSError:
            continue
    return False


def _resolve_path(path: str, default_ext: str = "") -> Path:
    """Resolve a user-supplied path.

    - If `path` is empty, place it in ~/jarvis_workspace with a timestamp name.
    - Relative paths land inside ~/jarvis_workspace.
    - Absolute paths are used as-is (after `_is_blocked`).
    - If `default_ext` is given and the path has no extension, add it.
    """
    if not path:
        stamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        return WORKSPACE / f"jarvis_{stamp}{default_ext}"
    p = Path(path).expanduser()
    if not p.is_absolute():
        p = WORKSPACE / p
    if default_ext and not p.suffix:
        p = p.with_suffix(default_ext)
    return p


def _write_text(p: Path, content: str) -> str:
    if _is_blocked(str(p)):
        return f"Refused: '{p}' is in a protected system directory."
    try:
        p.parent.mkdir(parents=True, exist_ok=True)
        p.write_text(content, encoding="utf-8")
        return f"Wrote {len(content)} chars to {p}"
    except Exception as exc:  # noqa: BLE001
        return f"Couldn't write {p}: {exc}"


# ── markdown-lite parser for PDF / DOCX / XLSX / PPTX ────
# Yields tuples of (kind, payload):
#   ("h1"/"h2"/"h3", str)
#   ("p", str)
#   ("ul", [str, ...])
#   ("quote", str)
#   ("table", [[str, ...], ...])   # first row is header
def _md_to_blocks(text: str):
    lines = text.splitlines()
    i = 0
    while i < len(lines):
        line = lines[i]
        stripped = line.strip()
        if not stripped:
            i += 1
            continue
        # Headings
        m = re.match(r"^(#{1,3})\s+(.*)$", stripped)
        if m:
            yield (f"h{len(m.group(1))}", m.group(2).strip())
            i += 1
            continue
        # Table — a row of `| col | col |` and a separator `| --- | --- |`
        if stripped.startswith("|") and i + 1 < len(lines) and re.match(
            r"^\s*\|?\s*(:?-{3,}:?\s*\|\s*)+:?-?\s*$", lines[i + 1]
        ):
            rows: list[list[str]] = []
            rows.append(_split_table_row(lines[i]))
            i += 2
            while i < len(lines) and lines[i].strip().startswith("|"):
                rows.append(_split_table_row(lines[i]))
                i += 1
            yield ("table", rows)
            continue
        # Blockquote
        if stripped.startswith(">"):
            buf = []
            while i < len(lines) and lines[i].strip().startswith(">"):
                buf.append(lines[i].strip().lstrip(">").lstrip())
                i += 1
            yield ("quote", " ".join(buf))
            continue
        # Unordered list
        if re.match(r"^[-*+]\s+", stripped):
            items: list[str] = []
            while i < len(lines) and re.match(r"^\s*[-*+]\s+", lines[i]):
                items.append(re.sub(r"^\s*[-*+]\s+", "", lines[i]).rstrip())
                i += 1
            yield ("ul", items)
            continue
        # Plain paragraph: gather until blank line / block-start
        buf = [stripped]
        i += 1
        while i < len(lines):
            nxt = lines[i].strip()
            if not nxt or nxt.startswith(("#", ">", "|", "-", "*", "+")):
                break
            buf.append(nxt)
            i += 1
        yield ("p", " ".join(buf))


def _split_table_row(line: str) -> list[str]:
    s = line.strip()
    if s.startswith("|"):
        s = s[1:]
    if s.endswith("|"):
        s = s[:-1]
    return [c.strip() for c in s.split("|")]


# ════════════════════════════════════════════════════════
#  Text-format tools (stdlib only)
# ════════════════════════════════════════════════════════

@tool(
    name="create_markdown",
    description=(
        "Create a Markdown file (.md). Pass the full Markdown content (headings "
        "with #, lists with -, etc.) and either a full path or a relative "
        "filename inside ~/jarvis_workspace."
    ),
    parameters={
        "type": "object",
        "properties": {
            "content": {"type": "string", "description": "Markdown body."},
            "path": {"type": "string", "description": "Output file path (.md)."},
            "title": {"type": "string", "description": "Optional H1 to prepend if content has no heading yet."},
        },
        "required": ["content"],
    },
)
def create_markdown(content: str, path: str = "", title: str = "") -> str:
    body = content
    if title and not re.search(r"^#{1,3}\s+", body.lstrip(), re.M):
        body = f"# {title}\n\n{body}"
    p = _resolve_path(path, ".md")
    return _write_text(p, body)


@tool(
    name="create_html",
    description=(
        "Create an HTML file (.html). If the content does not already contain "
        "<html>, it gets wrapped in a basic <!DOCTYPE html> document."
    ),
    parameters={
        "type": "object",
        "properties": {
            "content": {"type": "string", "description": "HTML body. Can be a full document or just the body fragment."},
            "path": {"type": "string", "description": "Output file path (.html)."},
            "title": {"type": "string", "description": "Optional <title> for the page."},
        },
        "required": ["content"],
    },
)
def create_html(content: str, path: str = "", title: str = "") -> str:
    if "<html" in content.lower():
        html = content
    else:
        title_tag = f"<title>{_html_escape(title)}</title>" if title else ""
        html = (
            "<!DOCTYPE html>\n"
            "<html lang=\"en\">\n"
            "<head>\n"
            "<meta charset=\"utf-8\">\n"
            f"{title_tag}\n"
            "</head>\n"
            "<body>\n"
            f"{content}\n"
            "</body>\n"
            "</html>\n"
        )
    p = _resolve_path(path, ".html")
    return _write_text(p, html)


@tool(
    name="create_csv",
    description=(
        "Create a CSV file (.csv) from a list of rows. The first row is treated "
        "as the header unless `headers` is provided separately."
    ),
    parameters={
        "type": "object",
        "properties": {
            "rows": {
                "type": "array",
                "items": {"type": "array", "items": {"type": "string"}},
                "description": "List of rows. Each row is a list of cell strings.",
            },
            "path": {"type": "string", "description": "Output file path (.csv)."},
            "headers": {
                "type": "array",
                "items": {"type": "string"},
                "description": "Optional header row. If given, it's prepended to rows.",
            },
        },
        "required": ["rows"],
    },
)
def create_csv(rows: list[list[str]], path: str = "", headers: list[str] | None = None) -> str:
    p = _resolve_path(path, ".csv")
    if _is_blocked(str(p)):
        return f"Refused: '{p}' is in a protected system directory."
    try:
        p.parent.mkdir(parents=True, exist_ok=True)
        with p.open("w", encoding="utf-8", newline="") as f:
            w = _csv.writer(f)
            if headers:
                w.writerow(headers)
            for r in rows:
                w.writerow(r)
        return f"Wrote {len(rows)} row(s) to {p}"
    except Exception as exc:  # noqa: BLE001
        return f"Couldn't write {p}: {exc}"


@tool(
    name="create_json",
    description=(
        "Create a JSON file (.json) from any JSON-serializable object (dict, "
        "list, string, number, etc.). Pretty-printed with indent=2."
    ),
    parameters={
        "type": "object",
        "properties": {
            "data": {"description": "Any JSON-serializable value (object, array, string, number)."},
            "path": {"type": "string", "description": "Output file path (.json)."},
        },
        "required": ["data"],
    },
)
def create_json(data, path: str = "") -> str:  # noqa: ANN001 — JSON shape is arbitrary
    p = _resolve_path(path, ".json")
    try:
        text = json.dumps(data, indent=2, ensure_ascii=False)
    except (TypeError, ValueError) as exc:
        return f"Data is not JSON-serializable: {exc}"
    return _write_text(p, text)


@tool(
    name="create_yaml",
    description=(
        "Create a YAML file (.yaml or .yml) from any Python object (dict, "
        "list, string, number, bool)."
    ),
    parameters={
        "type": "object",
        "properties": {
            "data": {"description": "Any YAML-serializable value."},
            "path": {"type": "string", "description": "Output file path (.yaml/.yml)."},
        },
        "required": ["data"],
    },
)
def create_yaml(data, path: str = "") -> str:  # noqa: ANN001
    try:
        import yaml  # type: ignore
    except ImportError:
        return "PyYAML is not installed. Run: pip install PyYAML"
    p = _resolve_path(path, ".yaml")
    if _is_blocked(str(p)):
        return f"Refused: '{p}' is in a protected system directory."
    try:
        p.parent.mkdir(parents=True, exist_ok=True)
        text = yaml.safe_dump(data, allow_unicode=True, sort_keys=False)
        p.write_text(text, encoding="utf-8")
        return f"Wrote {len(text)} chars to {p}"
    except Exception as exc:  # noqa: BLE001
        return f"Couldn't write {p}: {exc}"


@tool(
    name="create_xml",
    description=(
        "Create an XML file (.xml). The content should be the inner XML (no "
        "declaration needed); it gets wrapped in a <document> root if not "
        "already a full document."
    ),
    parameters={
        "type": "object",
        "properties": {
            "content": {"type": "string", "description": "XML body or fragment."},
            "path": {"type": "string", "description": "Output file path (.xml)."},
            "root": {"type": "string", "description": "Optional root element name when wrapping."},
        },
        "required": ["content"],
    },
)
def create_xml(content: str, path: str = "", root: str = "document") -> str:
    # Heuristic: if the content looks like XML (any open/close tags), treat
    # it as an XML fragment and just add a declaration if missing. If it's
    # plain text, wrap it in <root>…</root> with escaping.
    has_decl = content.lstrip().startswith("<?xml")
    looks_like_xml = bool(re.search(r"<[A-Za-z_][\w.-]*>", content)) and "</" in content
    if has_decl or looks_like_xml:
        if has_decl:
            xml = content
        else:
            xml = f"<?xml version=\"1.0\" encoding=\"utf-8\"?>\n{content.strip()}\n"
    else:
        xml = (
            f"<?xml version=\"1.0\" encoding=\"utf-8\"?>\n"
            f"<{root}>\n{_xml_escape(content)}\n</{root}>\n"
        )
    p = _resolve_path(path, ".xml")
    return _write_text(p, xml)


@tool(
    name="create_txt",
    description="Create a plain text file (.txt) with the given content.",
    parameters={
        "type": "object",
        "properties": {
            "content": {"type": "string"},
            "path": {"type": "string", "description": "Output file path (.txt)."},
        },
        "required": ["content"],
    },
)
def create_txt(content: str, path: str = "") -> str:
    p = _resolve_path(path, ".txt")
    return _write_text(p, content)


@tool(
    name="create_python_file",
    description=(
        "Create a Python source file (.py) with the given code. Optionally "
        "prepend a header line (e.g. shebang, encoding declaration)."
    ),
    parameters={
        "type": "object",
        "properties": {
            "content": {"type": "string", "description": "Python source code."},
            "path": {"type": "string", "description": "Output file path (.py)."},
            "header": {"type": "string", "description": "Optional first line, e.g. '#!/usr/bin/env python3'."},
        },
        "required": ["content"],
    },
)
def create_python_file(content: str, path: str = "", header: str = "") -> str:
    body = content
    if header:
        body = header.rstrip() + "\n" + body
    p = _resolve_path(path, ".py")
    return _write_text(p, body)


# Map of language -> file extension for create_code_file.
_LANG_EXT = {
    "python": "py", "py": "py",
    "javascript": "js", "js": "js", "node": "js",
    "typescript": "ts", "ts": "ts",
    "html": "html", "css": "css",
    "json": "json", "yaml": "yaml", "yml": "yaml",
    "java": "java", "kotlin": "kt", "swift": "swift",
    "c": "c", "cpp": "cpp", "c++": "cpp", "cc": "cpp",
    "csharp": "cs", "c#": "cs", "cs": "cs",
    "go": "go", "golang": "go",
    "rust": "rs", "rs": "rs",
    "ruby": "rb", "rb": "rb",
    "php": "php", "shell": "sh", "bash": "sh", "sh": "sh",
    "powershell": "ps1", "ps1": "ps1",
    "sql": "sql",
    "markdown": "md", "md": "md",
    "xml": "xml",
    "ini": "ini", "toml": "toml",
}


@tool(
    name="create_code_file",
    description=(
        "Create a source file in any common language. `language` is one of: "
        "python, javascript, typescript, html, css, java, kotlin, swift, c, "
        "cpp, csharp, go, rust, ruby, php, shell/bash, powershell, sql, "
        "markdown, xml, ini, toml, etc. The right extension is added if the "
        "path doesn't have one."
    ),
    parameters={
        "type": "object",
        "properties": {
            "content": {"type": "string", "description": "Source code."},
            "language": {"type": "string", "description": "Language name (e.g. 'python', 'js', 'go')."},
            "path": {"type": "string", "description": "Output file path (extension optional)."},
        },
        "required": ["content", "language"],
    },
)
def create_code_file(content: str, language: str, path: str = "") -> str:
    ext = _LANG_EXT.get(language.lower())
    if not ext:
        return f"Unknown language '{language}'. Use one of: {', '.join(sorted(set(_LANG_EXT.values())))}"
    p = _resolve_path(path, f".{ext}")
    return _write_text(p, content)


# ════════════════════════════════════════════════════════
#  Binary document tools (PDF, DOCX, XLSX, PPTX)
# ════════════════════════════════════════════════════════

@tool(
    name="create_pdf",
    description=(
        "Create a PDF file (.pdf) from a Markdown-lite body: use `#`/`##`/`###` "
        "for headings, blank-separated lines for paragraphs, `-` for bullet "
        "items, and `>` for blockquotes. Tables in `| col | col |` form also "
        "work."
    ),
    parameters={
        "type": "object",
        "properties": {
            "content": {"type": "string", "description": "Markdown-lite body."},
            "path": {"type": "string", "description": "Output file path (.pdf)."},
            "title": {"type": "string", "description": "Optional document title (for the PDF metadata)."},
        },
        "required": ["content"],
    },
)
def create_pdf(content: str, path: str = "", title: str = "") -> str:
    try:
        from reportlab.lib.pagesizes import LETTER
        from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
        from reportlab.lib.units import inch
        from reportlab.platypus import (
            SimpleDocTemplate, Paragraph, Spacer, ListFlowable, ListItem, Table, TableStyle,
        )
        from reportlab.lib import colors
    except ImportError:
        return "reportlab is not installed. Run: pip install reportlab"
    p = _resolve_path(path, ".pdf")
    if _is_blocked(str(p)):
        return f"Refused: '{p}' is in a protected system directory."
    try:
        p.parent.mkdir(parents=True, exist_ok=True)
        doc = SimpleDocTemplate(
            str(p), pagesize=LETTER,
            leftMargin=0.9 * inch, rightMargin=0.9 * inch,
            topMargin=0.9 * inch, bottomMargin=0.9 * inch,
            title=title or p.stem,
        )
        styles = getSampleStyleSheet()
        story: list = []
        for kind, payload in _md_to_blocks(content):
            if kind in ("h1", "h2", "h3"):
                level = int(kind[1])
                style = styles[f"Heading{level}"]
                story.append(Paragraph(_html_escape(payload), style))
                story.append(Spacer(1, 0.12 * inch))
            elif kind == "p":
                story.append(Paragraph(_html_escape(payload), styles["BodyText"]))
                story.append(Spacer(1, 0.10 * inch))
            elif kind == "ul":
                items = [ListItem(Paragraph(_html_escape(it), styles["BodyText"]),
                                  leftIndent=12, value="•") for it in payload]
                story.append(ListFlowable(items, bulletType="bullet", start="•"))
                story.append(Spacer(1, 0.08 * inch))
            elif kind == "quote":
                qstyle = ParagraphStyle("quote", parent=styles["BodyText"],
                                        leftIndent=18, textColor=colors.grey)
                story.append(Paragraph(_html_escape(payload), qstyle))
                story.append(Spacer(1, 0.10 * inch))
            elif kind == "table":
                t = Table(payload, hAlign="LEFT")
                t.setStyle(TableStyle([
                    ("BACKGROUND", (0, 0), (-1, 0), colors.lightgrey),
                    ("FONTNAME", (0, 0), (-1, 0), "Helvetica-Bold"),
                    ("GRID", (0, 0), (-1, -1), 0.4, colors.grey),
                    ("VALIGN", (0, 0), (-1, -1), "TOP"),
                    ("FONTSIZE", (0, 0), (-1, -1), 9),
                ]))
                story.append(t)
                story.append(Spacer(1, 0.12 * inch))
        doc.build(story)
        size = p.stat().st_size
        return f"Wrote PDF ({size} bytes) to {p}"
    except Exception as exc:  # noqa: BLE001
        return f"Couldn't write PDF {p}: {exc}"


@tool(
    name="create_docx",
    description=(
        "Create a Word document (.docx) from a Markdown-lite body: `#`/`##`/`###` "
        "for headings, blank-separated lines for paragraphs, `-` for bullets, "
        "`>` for blockquotes."
    ),
    parameters={
        "type": "object",
        "properties": {
            "content": {"type": "string", "description": "Markdown-lite body."},
            "path": {"type": "string", "description": "Output file path (.docx)."},
            "title": {"type": "string", "description": "Optional document title."},
        },
        "required": ["content"],
    },
)
def create_docx(content: str, path: str = "", title: str = "") -> str:
    try:
        from docx import Document
    except ImportError:
        return "python-docx is not installed. Run: pip install python-docx"
    p = _resolve_path(path, ".docx")
    if _is_blocked(str(p)):
        return f"Refused: '{p}' is in a protected system directory."
    try:
        p.parent.mkdir(parents=True, exist_ok=True)
        doc = Document()
        if title:
            doc.core_properties.title = title
        for kind, payload in _md_to_blocks(content):
            if kind == "h1":
                doc.add_heading(payload, level=1)
            elif kind == "h2":
                doc.add_heading(payload, level=2)
            elif kind == "h3":
                doc.add_heading(payload, level=3)
            elif kind == "p":
                doc.add_paragraph(payload)
            elif kind == "ul":
                for it in payload:
                    doc.add_paragraph(it, style="List Bullet")
            elif kind == "quote":
                para = doc.add_paragraph(payload)
                para.paragraph_format.left_indent = 720  # 0.5"
            elif kind == "table":
                if not payload:
                    continue
                t = doc.add_table(rows=len(payload), cols=len(payload[0]))
                t.style = "Light Grid Accent 1"
                for r, row in enumerate(payload):
                    for c, val in enumerate(row):
                        t.cell(r, c).text = val
        doc.save(str(p))
        return f"Wrote DOCX ({p.stat().st_size} bytes) to {p}"
    except Exception as exc:  # noqa: BLE001
        return f"Couldn't write DOCX {p}: {exc}"


@tool(
    name="create_xlsx",
    description=(
        "Create an Excel workbook (.xlsx) from one or more Markdown tables. "
        "Each `| col | col |` table becomes a separate sheet (Sheet1, Sheet2, "
        "etc.). If the content has no tables, a single sheet is created from "
        "the lines as one column."
    ),
    parameters={
        "type": "object",
        "properties": {
            "content": {"type": "string", "description": "Markdown with one or more `| col | col |` tables."},
            "path": {"type": "string", "description": "Output file path (.xlsx)."},
            "sheet_name_prefix": {"type": "string", "description": "Optional prefix for sheet names. Default 'Sheet'."},
        },
        "required": ["content"],
    },
)
def create_xlsx(content: str, path: str = "", sheet_name_prefix: str = "Sheet") -> str:
    try:
        from openpyxl import Workbook
    except ImportError:
        return "openpyxl is not installed. Run: pip install openpyxl"
    p = _resolve_path(path, ".xlsx")
    if _is_blocked(str(p)):
        return f"Refused: '{p}' is in a protected system directory."
    try:
        p.parent.mkdir(parents=True, exist_ok=True)
        wb = Workbook()
        # Drop the default sheet — we'll add named ones.
        wb.remove(wb.active)

        tables = [payload for kind, payload in _md_to_blocks(content) if kind == "table"]
        if not tables:
            # Fallback: one sheet, one column, one row per non-empty line.
            wb.create_sheet(f"{sheet_name_prefix}1")
            ws = wb.active
            for r, line in enumerate(content.splitlines(), start=1):
                if line.strip():
                    ws.cell(row=r, column=1, value=line.strip())
        else:
            for idx, rows in enumerate(tables, start=1):
                ws = wb.create_sheet(f"{sheet_name_prefix}{idx}")
                for r, row in enumerate(rows, start=1):
                    for c, val in enumerate(row, start=1):
                        ws.cell(row=r, column=c, value=val)
        wb.save(str(p))
        return f"Wrote XLSX ({p.stat().st_size} bytes) to {p}"
    except Exception as exc:  # noqa: BLE001
        return f"Couldn't write XLSX {p}: {exc}"


@tool(
    name="create_pptx",
    description=(
        "Create a PowerPoint deck (.pptx) from a Markdown-lite body. Each "
        "`# Heading 1` starts a new slide; the paragraphs and bullets that "
        "follow it become the slide's content (title from the H1, body from "
        "the rest)."
    ),
    parameters={
        "type": "object",
        "properties": {
            "content": {"type": "string", "description": "Markdown-lite body."},
            "path": {"type": "string", "description": "Output file path (.pptx)."},
            "title": {"type": "string", "description": "Optional deck title (for the file metadata)."},
        },
        "required": ["content"],
    },
)
def create_pptx(content: str, path: str = "", title: str = "") -> str:
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
    except ImportError:
        return "python-pptx is not installed. Run: pip install python-pptx"
    p = _resolve_path(path, ".pptx")
    if _is_blocked(str(p)):
        return f"Refused: '{p}' is in a protected system directory."
    try:
        p.parent.mkdir(parents=True, exist_ok=True)
        prs = Presentation()
        if title:
            prs.core_properties.title = title
        blank = prs.slide_layouts[6]  # blank layout
        # Build slides from the markdown.
        blocks = list(_md_to_blocks(content))
        has_h1 = any(kind == "h1" for kind, _ in blocks)
        # Cover slide only when the caller gave an explicit title, or the
        # content has no H1 to open with — otherwise every deck would start
        # with a spurious "Presentation" slide in front of the real title.
        if title or not has_h1:
            title_slide = prs.slides.add_slide(blank)
            tb = title_slide.shapes.add_textbox(Inches(0.7), Inches(2.5), Inches(8.6), Inches(1.2))
            tf = tb.text_frame
            tf.text = title or "Presentation"
            for para in tf.paragraphs:
                for run in para.runs:
                    run.font.size = Pt(40)
                    run.font.bold = True

        current_title: str | None = None
        body_items: list[str] = []
        first_h1_seen = False

        def _flush_slide():
            nonlocal current_title, body_items
            if current_title is None and not body_items:
                return
            slide = prs.slides.add_slide(blank)
            tb = slide.shapes.add_textbox(Inches(0.7), Inches(0.6), Inches(8.6), Inches(1.0))
            tf = tb.text_frame
            tf.text = current_title or ""
            for para in tf.paragraphs:
                for run in para.runs:
                    run.font.size = Pt(28)
                    run.font.bold = True
            if body_items:
                bb = slide.shapes.add_textbox(Inches(0.7), Inches(1.8), Inches(8.6), Inches(5.0))
                btf = bb.text_frame
                btf.word_wrap = True
                for idx, it in enumerate(body_items):
                    if idx == 0:
                        btf.text = f"• {it}"
                    else:
                        p = btf.add_paragraph()
                        p.text = f"• {it}"
                    for para in btf.paragraphs:
                        for run in para.runs:
                            run.font.size = Pt(18)
            current_title = None
            body_items = []

        for kind, payload in blocks:
            if kind == "h1":
                if first_h1_seen:
                    _flush_slide()
                first_h1_seen = True
                current_title = payload
            elif kind in ("h2", "h3"):
                body_items.append(f"**{payload}**")
            elif kind == "ul":
                body_items.extend(payload)
            elif kind == "p":
                body_items.append(payload)
            elif kind == "quote":
                body_items.append(f"“{payload}”")
        _flush_slide()
        prs.save(str(p))
        return f"Wrote PPTX ({p.stat().st_size} bytes) to {p}"
    except Exception as exc:  # noqa: BLE001
        return f"Couldn't write PPTX {p}: {exc}"


# ════════════════════════════════════════════════════════
#  Image generation (gemini-2.5-flash-image)
# ════════════════════════════════════════════════════════

class _ImageClient:
    """Thin POST helper with the same key-rotation as Brain._post_stream."""

    def __init__(self) -> None:
        self.session = requests.Session()
        self._key_index = 0

    def _mark_dead(self) -> None:
        keys = CONFIG.all_keys()
        if len(keys) <= 1:
            return
        self._key_index = (self._key_index + 1) % len(keys)

    def generate(self, prompt: str, aspect_ratio: str = "1:1") -> bytes | None:
        keys = CONFIG.all_keys()
        if not keys:
            return None
        last_err = ""
        # Only image-capable models here. The old fallback (gemini-2.0-flash)
        # is a *text* model and never returns image bytes, so it made the
        # "fallback" a guaranteed no-op — dropped.
        for model in ["gemini-2.5-flash-image"]:
            for ki, key in enumerate(keys):
                if ki < self._key_index:
                    continue
                url = f"{GEMINI_BASE}/{model}:generateContent"
                try:
                    r = self.session.post(
                        url, params={"key": key}, timeout=120,
                        json={
                            "contents": [{"role": "user", "parts": [{"text": prompt}]}],
                            "generationConfig": {
                                "responseModalities": ["IMAGE", "TEXT"],
                                # Honor the requested aspect ratio instead of
                                # silently dropping it (was accepted but unused).
                                "imageConfig": {"aspectRatio": aspect_ratio},
                            },
                        },
                    )
                except requests.RequestException as exc:
                    last_err = f"network {exc}"
                    continue
                if r.status_code == 200:
                    self._key_index = ki
                    return _extract_png(r.json())
                if r.status_code in (429, 500, 502, 503, 504):
                    last_err = f"{r.status_code} {r.text[:120]}"
                    self._mark_dead()
                    continue
                last_err = f"{r.status_code} {r.text[:120]}"
                break
        return None  # caller logs last_err via raise if needed


def _extract_png(payload: dict) -> bytes | None:
    try:
        for part in (payload.get("candidates", [{}])[0]
                     .get("content", {}).get("parts", []) or []):
            inline = part.get("inlineData") or part.get("inline_data")
            if inline and inline.get("data"):
                return base64.b64decode(inline["data"])
    except (IndexError, KeyError, TypeError):
        pass
    return None


def _safe_filename(s: str, maxlen: int = 40) -> str:
    s = re.sub(r"[^A-Za-z0-9_.-]+", "_", s).strip("._")
    return (s or "image")[:maxlen]


@tool(
    name="generate_image",
    description=(
        "Generate an image from a text prompt using Gemini 2.5 Flash Image. "
        "Saves a PNG to ~/jarvis_workspace/images/ unless `path` is given. "
        "Returns the absolute path."
    ),
    parameters={
        "type": "object",
        "properties": {
            "prompt": {"type": "string", "description": "A clear description of the image to create."},
            "path": {"type": "string", "description": "Output file path (.png). Optional."},
            "aspect_ratio": {"type": "string", "description": "1:1, 16:9, 9:16, 4:3, 3:4. Default 1:1.", "default": "1:1"},
        },
        "required": ["prompt"],
    },
)
def generate_image(prompt: str, path: str = "", aspect_ratio: str = "1:1") -> str:
    if not CONFIG.has_key and not CONFIG.has_backup_key():
        return "No Gemini API key configured. Open Settings and paste one."
    out = _resolve_path(path, ".png")
    if out.suffix.lower() != ".png":
        out = out.with_suffix(".png")
    if _is_blocked(str(out)):
        return f"Refused: '{out}' is in a protected system directory."
    try:
        out.parent.mkdir(parents=True, exist_ok=True)
        client = _ImageClient()
        try:
            png = client.generate(prompt, aspect_ratio)
        finally:
            client.session.close()  # don't leak the connection pool per call
        if not png:
            return "Image generation failed — the API returned no image bytes (rate-limited or content blocked)."
        out.write_bytes(png)
        return f"Image saved to {out} ({len(png)} bytes)"
    except Exception as exc:  # noqa: BLE001
        return f"Couldn't generate image: {exc}"


# ════════════════════════════════════════════════════════
#  Code execution (sandboxed, approval-gated)
# ════════════════════════════════════════════════════════

@tool(
    name="run_python",
    description=(
        "Execute a Python code snippet and return its stdout/stderr. Runs in "
        "a temp directory with a 10-second timeout. Requires user approval. "
        "Use for one-off calculations, quick data wrangling, or to test code "
        "before saving it to a file."
    ),
    requires_approval=True,
    destructive=True,
    parameters={
        "type": "object",
        "properties": {
            "code": {"type": "string", "description": "Python source code to execute."},
        },
        "required": ["code"],
    },
)
def run_python(code: str) -> str:
    with tempfile.TemporaryDirectory(prefix="jarvis_py_") as tmp:
        try:
            proc = subprocess.run(
                [sys.executable, "-I", "-c", code],
                cwd=tmp, capture_output=True, text=True,
                timeout=10, shell=False,
            )
        except subprocess.TimeoutExpired:
            return "Code timed out after 10s."
        except Exception as exc:  # noqa: BLE001
            return f"Couldn't run code: {exc}"
        out = (proc.stdout or "") + (proc.stderr or "")
        if not out.strip():
            return f"(exit {proc.returncode}, no output)"
        return out.strip()[:4000] + (
            f"\n(exit {proc.returncode})" if proc.returncode != 0 else ""
        )
